"""Shape analysis for PPTX watermark detection."""

from pptx.enum.shapes import MSO_SHAPE_TYPE
from .utils import emu_to_inches, get_shape_position_percentage, is_bottom_right_corner


def analyze_hyperlinks(shape):
    """Extract hyperlinks from a shape."""
    hyperlinks = []

    # Check if the shape itself has a click action (hyperlink)
    if hasattr(shape, "click_action") and shape.click_action:
        action = shape.click_action
        if hasattr(action, "hyperlink") and action.hyperlink:
            hyperlinks.append(f"Click action: {action.hyperlink.address}")

    # Check for hyperlinks in text frames
    if hasattr(shape, "text_frame"):
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if hasattr(run, "hyperlink") and run.hyperlink:
                        if run.hyperlink.address:
                            hyperlinks.append(
                                f"Text hyperlink: {run.hyperlink.address}"
                            )
        except Exception:
            pass

    return hyperlinks


def analyze_shape(shape, slide_width, slide_height, indent=0):
    """Analyze a single shape and return details."""
    prefix = "  " * indent
    result = []

    # Get shape type
    shape_type = str(shape.shape_type) if hasattr(shape, "shape_type") else "Unknown"

    # Get position
    left = shape.left if shape.left is not None else 0
    top = shape.top if shape.top is not None else 0
    width = shape.width if shape.width is not None else 0
    height = shape.height if shape.height is not None else 0

    left_pct, top_pct, right_pct, bottom_pct = get_shape_position_percentage(
        shape, slide_width, slide_height
    )

    result.append(f"{prefix}Shape: {shape.name}")
    result.append(f"{prefix}  Type: {shape_type}")
    result.append(
        f"{prefix}  Position (EMUs): left={left}, top={top}, width={width}, height={height}"
    )
    result.append(
        f'{prefix}  Position (inches): left={emu_to_inches(left):.2f}", '
        f'top={emu_to_inches(top):.2f}", width={emu_to_inches(width):.2f}", '
        f'height={emu_to_inches(height):.2f}"'
    )

    if left_pct is not None:
        result.append(
            f"{prefix}  Position (%): left={left_pct:.1f}%, top={top_pct:.1f}%, "
            f"right={right_pct:.1f}%, bottom={bottom_pct:.1f}%"
        )

    # Check if in bottom-right corner
    if is_bottom_right_corner(shape, slide_width, slide_height, 70):
        result.append(f"{prefix}  *** BOTTOM-RIGHT CORNER (>70%) ***")

    # Get text content
    if hasattr(shape, "text") and shape.text:
        text = shape.text[:100] + "..." if len(shape.text) > 100 else shape.text
        result.append(f"{prefix}  Text: '{text}'")
        if "gamma" in shape.text.lower():
            result.append(f"{prefix}  *** CONTAINS 'GAMMA' IN TEXT ***")

    # Get hyperlinks
    hyperlinks = analyze_hyperlinks(shape)
    if hyperlinks:
        for link in hyperlinks:
            result.append(f"{prefix}  Hyperlink: {link}")
            if "gamma" in link.lower():
                result.append(f"{prefix}  *** CONTAINS 'GAMMA' IN HYPERLINK ***")

    # Check for image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        result.append(f"{prefix}  *** THIS IS AN IMAGE ***")
        try:
            if hasattr(shape, "image"):
                result.append(f"{prefix}  Image format: {shape.image.content_type}")
                result.append(f"{prefix}  Image size: {len(shape.image.blob)} bytes")
        except Exception as e:
            result.append(f"{prefix}  Image info error: {e}")

    # Handle group shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        result.append(f"{prefix}  *** THIS IS A GROUP SHAPE ***")
        for subshape in shape.shapes:
            result.extend(
                analyze_shape(subshape, slide_width, slide_height, indent + 1)
            )

    return result
